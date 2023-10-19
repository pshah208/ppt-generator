import streamlit as st
import base64
import openai
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv
from langchain.document_loaders import DirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.vectorstores import Chroma
from langchain.vectorstores import FAISS
from langchain.embeddings import OpenAIEmbeddings
from langchain.chat_models import ChatOpenAI
from langchain.chains.summarize import load_summarize_chain
from langchain.chains import RetrievalQA

load_dotenv()
# Get an OpenAI API Key before continuing
if "openai_api_key" in st.secrets:
    openai.api_key = st.secrets.openai_api_key
else:
    openai.api_key = st.sidebar.text_input("OpenAI API Key", type="password")
if not openai.api_key:
    st.info("Enter an OpenAI API Key to continue")
    st.stop()

st.title("PPT Generator - Acharya")

# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

# Load documents from local directory
loader = DirectoryLoader('./doc/', glob="**/[!.]*")
docs = loader.load()


splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=20)

documents = splitter.split_documents(docs)


 # Create vector embeddings and store them in a vector database
vectorstore = FAISS.from_documents(documents, embedding=OpenAIEmbeddings(openai_api_key=openai.api_key))                                   
   
topic = st.text_input("Enter the topic for your presentation:")
retriever = vectorstore.as_retriever(k=1)
qa = RetrievalQA.from_chain_type(llm=ChatOpenAI(openai_api_key=openai.api_key), chain_type="stuff", retriever=retriever)

def generate_slide_titles(topic, qa):
    prompt = f"Generate 5 slide titles for the topic '{topic}' and retrieve slide titles with the help of QA: '{qa}'."
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        max_tokens=200,
    )
    return response['choices'][0]['text'].split("\n")

def generate_slide_content(slide_title, qa):
    prompt = f"Generate content for the slide: '{slide_title}' and retrieve information from documents: '{qa}'."
    response = openai.Completion.create(
        model="text-davinci-003",
        prompt=prompt,
        max_tokens=500,  # Adjust as needed based on the desired content length
    )
    return response['choices'][0]['text']


def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    prs.save(f"doc/{topic}_presentation.pptx")
    

def main():
    
    
    
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic, qa)
        filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
        print("Slide Title: ", filtered_slide_titles)
        slide_contents = [generate_slide_content(title, qa) for title in filtered_slide_titles]
        print("Slide Contents: ", slide_contents)
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation generated successfully!")

        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    ppt_filename = f"doc/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


if __name__ == "__main__":
    main()



