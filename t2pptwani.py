import streamlit as st
import base64
import openai
import pptx
from pptx.util import Inches, Pt
import os
from dotenv import load_dotenv
from langchain.document_loaders import DirectoryLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.embeddings import OpenAIEmbeddings
from langchain.chat_models import ChatOpenAI
from langchain.utilities.dalle_image_generator import DallEAPIWrapper

load_dotenv()
# Get an OpenAI API Key before continuing
if "openai.api_key" in st.secrets:
    openai.api_key = st.secrets.openai_api_key
else:
    openai.api_key = st.sidebar.text_input("OpenAI API Key", type="password")
if not openai.api_key:
    st.title("PPT Generator - Acharya")
    st.info("Enter an OpenAI API Key to continue")
    st.info("If you are not sure on how to get your OpenAI API key:")
    st.info( " 1) Please visit https://platform.openai.com/account/api-keys")
    st.info(" 2) Click on 'Create new key' and copy and save the key in a safe location")
    st.stop()

st.title("PPT Generator - Acharya")
topic = st.text_input("Enter the topic for your presentation:")
# Define custom formatting options
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)


llm = ChatOpenAI(openai_api_key=openai.api_key, max_tokens = 500)

def generate_slide_titles(topic):
   
  prompt_template = """
  Generate 5 engaging slide titles for a powerpoint presentation about {topic}
  """
  
  chain = LLMChain(llm=llm, prompt=PromptTemplate.from_template(prompt_template))

  result = chain.run(topic=topic)

  # Parse the output into a list
  titles = result.split("\n")

  return titles

def generate_slide_content(slide_title):
    text_prompt = f"Generate content for the slide: '{slide_title}' ."
    text_response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=text_prompt,
        max_tokens=500,  # Adjust as needed based on the desired content length
    )
    content_text = text_response['choices'][0]['text']
  
    return content_text
  

def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()
    slide_layout = prs.slide_layouts[1]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic

    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content
       # image = slide.shapes.add_picture(image_url, 0, 0)

        # Customize font size for titles and content
        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE

    prs.save(f"doc/{topic}_presentation.pptx")
    

def main():
    
    generate_button = st.button("Generate Presentation")

    if generate_button and topic:
        st.info("Generating presentation... Please wait.")
        slide_titles = generate_slide_titles(topic)
        filtered_slide_titles= [item for item in slide_titles if item.strip() != '']
        print("Slide Title: ", filtered_slide_titles)
        slide_contents = [generate_slide_content(title) for title in filtered_slide_titles]
      #  images = [generate_image(title) for title in filtered_slide_titles]
        print("Slide Contents: ", slide_contents)
        create_presentation(topic, filtered_slide_titles, slide_contents)
        print("Presentation generated successfully!")

        st.success("Presentation generated successfully!")
        st.markdown(get_ppt_download_link(topic), unsafe_allow_html=True)

def get_ppt_download_link(topic):
    ppt_filename = f"doc/{topic}_presentation.pptx"

    with open(ppt_filename, "rb") as file:
        ppt_contents = file.read()

    b64_ppt = base64.b64encode(ppt_contents).decode()
    return f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_ppt}" download="{ppt_filename}">Download the PowerPoint Presentation</a>'


if __name__ == "__main__":
    main()
